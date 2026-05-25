#!/usr/bin/env python3
"""
Knowledge Base ingestion pipeline.

Supports PDF, DOCX, XLSX, PPTX, TXT, MD, source code, and legacy
Office formats. Extracts text, chunks with overlap, generates embeddings,
and upserts to Qdrant. Tracks file state in the IngestRegistry.

Usage:
    python ingest.py --docs /path/to/docs
    python ingest.py --docs /path --product my-product --workers 4
    python ingest.py --file /path/to/file.pdf
    python ingest.py --clean  # clear KB before re-ingesting
"""

import argparse
import asyncio
import logging
import os
import sys
import time
import uuid
from pathlib import Path

# ── CRÍTICO: carrega o .env ANTES de qualquer import que leia variáveis ──────
_project_root = Path(__file__).parent.parent
sys.path.insert(0, str(_project_root))
from config.bootstrap_env import bootstrap_env
bootstrap_env()

# Agora adiciona o server/ ao path para imports de embed_client e vector_store
sys.path.insert(0, str(_project_root / "server"))

log = logging.getLogger("kb-ingest")
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    handlers=[logging.StreamHandler(sys.stdout)],
)


# ── Extensões suportadas
# ──────────────────────────────────────────────────────
EXT_TYPE_MAP = {
    ".pdf": "pdf",
    ".docx": "docx",
    ".doc": "doc",        # legacy Word 97-2003
    ".xlsx": "xlsx",
    ".xls": "xls",        # legacy Excel 97-2003
    ".pptx": "pptx",
    ".ppt": "ppt",        # legacy PowerPoint 97-2003
    ".txt": "txt",
    ".md": "txt",
    ".rst": "txt",
    ".py": "code",
    ".ts": "code",
    ".js": "code",
    ".java": "code",
    ".go": "code",
    ".rs": "code",
    ".cpp": "code",
    ".c": "code",
    ".cs": "code",
    ".yaml": "code",
    ".yml": "code",
    ".json": "code",
    ".xml": "code",
    ".sh": "code",
    ".sql": "code",
    ".odt": "odt",        # OpenDocument Text
    ".ods": "ods",        # OpenDocument Spreadsheet
    ".odp": "odp",        # OpenDocument Presentation
    ".wpd": "wpd",        # WordPerfect
    ".zip": "zip",        # ZIP archive
}

# Chunk settings por tipo
_DEFAULT_CHUNK_SIZE = int(os.getenv("INGEST_CHUNK_SIZE_DEFAULT", "600"))
_DEFAULT_CHUNK_OVERLAP = int(os.getenv("INGEST_CHUNK_OVERLAP_DEFAULT", "80"))

CHUNK_SETTINGS = {
    "pdf":  {"size": int(os.getenv("INGEST_CHUNK_SIZE_PDF",  "800")),
             "overlap": int(os.getenv("INGEST_CHUNK_OVERLAP_PDF", "100"))},
    "docx": {"size": int(os.getenv("INGEST_CHUNK_SIZE_DOCX", "700")),
             "overlap": int(os.getenv("INGEST_CHUNK_OVERLAP_DOCX", "80"))},
    "xlsx": {"size": int(os.getenv("INGEST_CHUNK_SIZE_XLSX", "500")),
             "overlap": int(os.getenv("INGEST_CHUNK_OVERLAP_XLSX", "50"))},
    "pptx": {"size": int(os.getenv("INGEST_CHUNK_SIZE_PPTX", "600")),
             "overlap": int(os.getenv("INGEST_CHUNK_OVERLAP_PPTX", "80"))},
    "txt":  {"size": int(os.getenv("INGEST_CHUNK_SIZE_TXT",  "600")),
             "overlap": int(os.getenv("INGEST_CHUNK_OVERLAP_TXT",  "80"))},
    "code": {"size": int(os.getenv("INGEST_CHUNK_SIZE_CODE", "400")),
             "overlap": int(os.getenv("INGEST_CHUNK_OVERLAP_CODE", "60"))},
}


# ─────────────────────────────────────────────────────────────────────────────
# EXTRATORES
# ─────────────────────────────────────────────────────────────────────────────


def extract_pdf(path: Path) -> list[dict]:
    """Extract text from a PDF file preserving page structure.

    Attempts docling first (best for complex PDFs with tables/figures),
    then falls back to PyMuPDF (fitz) for simpler documents.

    Args:
        path: Path to the PDF file.

    Returns:
        List of dicts with 'text' and 'page' keys, or empty list on error.
    """
    chunks_raw = []
    try:
        # Tenta docling primeiro (melhor para PDFs complexos)
        from docling.document_converter import DocumentConverter

        converter = DocumentConverter()
        result = converter.convert(str(path))
        text = result.document.export_to_markdown()
        chunks_raw.append({"text": text, "page": None})
        log.info(f"  docling OK: {path.name}")
    except ImportError:
        pass
    except Exception as e:
        log.warning(f"  docling falhou ({e}), tentando PyMuPDF")

    if not chunks_raw:
        try:
            import fitz  # PyMuPDF

            doc = fitz.open(str(path))
            for page_num, page in enumerate(doc, 1):
                text = page.get_text("text").strip()
                if text:
                    chunks_raw.append({"text": text, "page": page_num})
            log.info(f"  PyMuPDF OK: {path.name} ({len(chunks_raw)} páginas)")
        except ImportError:
            log.error(
                "  Instale docling ou pymupdf: pip install docling pymupdf"
            )
            return []
        except Exception as e:
            log.error(f"  Erro ao extrair PDF: {e}")
            return []

    return chunks_raw


def extract_docx(path: Path) -> list[dict]:
    """Extract text from a DOCX file using python-docx.

    Args:
        path: Path to the DOCX file.

    Returns:
        List with a single dict containing concatenated paragraph text,
        or empty list on error.
    """
    try:
        from docx import Document

        doc = Document(str(path))
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        return [{"text": "\n".join(paragraphs), "page": None}]
    except ImportError:
        log.error("  Instale python-docx: pip install python-docx")
        return []
    except Exception as e:
        log.error(f"  Erro ao extrair DOCX: {e}")
        return []


def extract_xlsx(path: Path) -> list[dict]:
    """Extract text from an XLSX file using openpyxl.

    Reads all sheets in read-only mode, concatenating rows with tab
    separators.

    Args:
        path: Path to the XLSX file.

    Returns:
        List of dicts with 'text' and 'page' (sheet name) keys,
        or empty list on error.
    """
    try:
        import openpyxl

        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        results = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                row_text = "\t".join(
                    str(c) if c is not None else "" for c in row
                )
                if row_text.strip():
                    rows.append(row_text)
            if rows:
                results.append(
                    {
                        "text": (
                            f"[Planilha: {sheet_name}]\n" + "\n".join(rows)
                        ),
                        "page": sheet_name,
                    }
                )
        return results
    except ImportError:
        log.error("  Instale openpyxl: pip install openpyxl")
        return []
    except Exception as e:
        log.error(f"  Erro ao extrair XLSX: {e}")
        return []


def extract_pptx(path: Path) -> list[dict]:
    """Extract text from a PPTX file using python-pptx.

    Iterates through slides collecting text from all shapes.

    Args:
        path: Path to the PPTX file.

    Returns:
        List of dicts with 'text' and 'page' (slide number) keys,
        or empty list on error.
    """
    try:
        from pptx import Presentation

        prs = Presentation(str(path))
        results = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                results.append({"text": "\n".join(texts), "page": i})
        return results
    except ImportError:
        log.error("  Instale python-pptx: pip install python-pptx")
        return []
    except Exception as e:
        log.error(f"  Erro ao extrair PPTX: {e}")
        return []


def extract_text(path: Path) -> list[dict]:
    """Extract text from a plain text file.

    Args:
        path: Path to the text file.

    Returns:
        List with a single dict containing the full file text,
        or empty list on error.
    """
    try:
        text = path.read_text(encoding="utf-8", errors="replace")
        return [{"text": text, "page": None}]
    except Exception as e:
        log.error(f"  Erro ao ler texto: {e}")
        return []


def extract_code(path: Path) -> list[dict]:
    """Extract source code wrapped in a language-marked code block.

    Adds a language identifier (e.g., ```python) for better context.

    Args:
        path: Path to the source code file.

    Returns:
        List with a single dict containing the code in a fenced block,
        or empty list on error.
    """
    try:
        text = path.read_text(encoding="utf-8", errors="replace")
        lang = path.suffix.lstrip(".")
        return [{"text": f"```{lang}\n{text}\n```", "page": None}]
    except Exception as e:
        log.error(f"  Erro ao ler código: {e}")
        return []


from ingest.parsers.legacy_office import (
    extract_doc, extract_xls, extract_ppt,
    extract_odt, extract_ods, extract_odp, extract_wpd,
)
from ingest.parsers.zip_handler import extract_zip

EXTRACTORS = {
    "pdf": extract_pdf,
    "docx": extract_docx,
    "doc": extract_doc,        # legacy Word 97-2003
    "xlsx": extract_xlsx,
    "xls": extract_xls,        # legacy Excel 97-2003
    "pptx": extract_pptx,
    "ppt": extract_ppt,        # legacy PowerPoint 97-2003
    "txt": extract_text,
    "code": extract_code,
    "odt": extract_odt,        # OpenDocument Text
    "ods": extract_ods,        # OpenDocument Spreadsheet
    "odp": extract_odp,        # OpenDocument Presentation
    "wpd": extract_wpd,        # WordPerfect
    "zip": extract_zip,        # ZIP archive
}


# ─────────────────────────────────────────────────────────────────────────────
# CHUNKING
# ─────────────────────────────────────────────────────────────────────────────


def chunk_text(text: str, file_type: str) -> list[str]:
    """Divide texto em chunks com overlap usando LangChain."""
    settings = CHUNK_SETTINGS.get(file_type, {"size": 600, "overlap": 80})
    try:
        from langchain_text_splitters import RecursiveCharacterTextSplitter

        splitter = RecursiveCharacterTextSplitter(
            chunk_size=settings["size"],
            chunk_overlap=settings["overlap"],
            separators=["\n\n", "\n", ". ", " ", ""],
        )
        return splitter.split_text(text)
    except ImportError:
        # Fallback manual se langchain não estiver instalado
        size = settings["size"]
        overlap = settings["overlap"]
        chunks = []
        start = 0
        while start < len(text):
            end = start + size
            chunks.append(text[start:end])
            start = end - overlap
        return chunks


# ─────────────────────────────────────────────────────────────────────────────
# DOCUMENT PARSING HELPER
# ─────────────────────────────────────────────────────────────────────────────


def parse_document(file_path: Path) -> list[dict]:
    """
    Parse a document file into a list of chunk dicts.

    Returns a list of dicts with keys: text, file_type, page (optional).
    Returns an empty list if the file type is unsupported.
    """
    ext = file_path.suffix.lower()
    file_type = EXT_TYPE_MAP.get(ext)
    if not file_type:
        return []
    extractor = EXTRACTORS.get(file_type)
    if not extractor:
        return []
    return extractor(file_path)


# ─────────────────────────────────────────────────────────────────────────────
# PIPELINE PRINCIPAL
# ─────────────────────────────────────────────────────────────────────────────


async def process_file(
    file_path: Path,
    docs_root: Path,
    store,
    registry,
    product_override: str | None = None,
    force: bool = False,
) -> tuple[int, str]:
    """Process a single file: extract, chunk, embed, and upsert to Qdrant.

    Classifies the file (product, doc_type), computes embeddings via batch
    API, upserts chunks to the vector store, and updates the registry.

    Args:
        file_path: Path to the file to process.
        docs_root: Root directory for relative path computation.
        store: VectorStore instance for upsert operations.
        registry: IngestRegistry for tracking file state.
        product_override: Optional product name override.
        force: If True, re-ingest even if no changes detected.

    Returns:
        Tuple of (chunk_count, status) where status is 'ok', 'skipped',
        or 'error'.
    """
    from ingest.classifier import classify
    from kb_server.embed_client import get_embeddings_batch

    ext = file_path.suffix.lower()
    file_type = EXT_TYPE_MAP.get(ext)
    if not file_type:
        return 0, "skipped"

    meta = classify(file_path, docs_root, product_override)
    product = meta["product"]
    doc_type = meta["doc_type"]
    version = meta.get("version")  # FASE 13: Optional version field
    vendor = meta.get("vendor", "")
    subsystem = meta.get("subsystem", "")
    source_file = str(file_path.relative_to(docs_root))

    # ── Verifica se precisa ingerir ─
    # ───────────────────────────────────────────
    if not force:
        needs, reason = registry.needs_ingest(file_path, source_file)
        if not needs:
            log.debug(f"  SKIP: {source_file} ({reason})")
            return 0, "skipped"
        vendor_tag = f"[{vendor:10}] " if vendor else ""
        log.info(f"{vendor_tag}[{doc_type:16}] [{product:20}] {source_file}  ({reason})")
    else:
        log.info(f"[{doc_type:16}] [{product:20}] {source_file}  (forçado)")

    try:
        extractor = EXTRACTORS[file_type]
        raw_sections = extractor(file_path)
        if not raw_sections:
            log.warning(f"  Sem conteúdo extraído: {source_file}")
            registry.mark_error(
                file_path,
                source_file,
                "sem conteúdo extraído",
                file_type,
                product,
                doc_type,
            )
            return 0, "error"

        # Remove chunks anteriores do mesmo arquivo
        await store.delete_document(source_file)

        # Gera chunks
        all_chunks_data = []
        chunk_index = 0
        for section in raw_sections:
            text_chunks = chunk_text(section["text"], file_type)
            for chunk_text_content in text_chunks:
                if len(chunk_text_content.strip()) < 30:
                    continue
                chunk_data = {
                    "chunk_id": str(uuid.uuid4()),
                    "text": chunk_text_content,
                    "source_file": source_file,
                    "file_type": file_type,
                    "product": product,
                    "doc_type": doc_type,
                    "vendor": vendor,
                    "subsystem": subsystem,
                    "page": section.get("page"),
                    "chunk_index": chunk_index,
                }
                # FASE 13: Add version if available
                if version:
                    chunk_data["version"] = version
                all_chunks_data.append(chunk_data)
                chunk_index += 1

        if not all_chunks_data:
            registry.mark_error(
                file_path,
                source_file,
                "chunks vazios após split",
                file_type,
                product,
                doc_type,
            )
            return 0, "error"

        texts = [c["text"] for c in all_chunks_data]
        vectors = await get_embeddings_batch(texts, batch_size=32)
        for chunk_data, vector in zip(all_chunks_data, vectors):
            chunk_data["vector"] = vector

        await store.upsert_chunks(all_chunks_data)
        registry.mark_ok(
            file_path,
            source_file,
            len(all_chunks_data),
            file_type,
            product,
            doc_type,
        )
        log.info(f"  ✓ {len(all_chunks_data)} chunks")
        return len(all_chunks_data), "ok"

    except Exception as e:
        log.error(f"  ✗ {source_file}: {e}", exc_info=True)
        registry.mark_error(
            file_path, source_file, str(e), file_type, product, doc_type
        )
        return 0, "error"


async def run_ingest(
    docs_path: Path,
    product: str | None = None,
    workers: int = 2,
    single_file: Path | None = None,
    clean: bool = False,
    force: bool = False,
    sync: bool = False,
):
    """Run the full ingest pipeline or process a single file.

    Connects to Qdrant and the registry, discovers files to process,
    runs them with bounded parallelism, logs a summary, and handles
    sync-deletion of removed files.

    Args:
        docs_path: Root directory containing documents to ingest.
        product: Optional product name override for all files.
        workers: Number of parallel file processors (default 2).
        single_file: Process a single file instead of a directory.
        clean: Clear KB and registry before re-ingesting.
        force: Re-ingest even if no changes detected.
        sync: Mark deleted files removed from disk in registry.
    """
    from ingest.core.metadata import IngestRegistry
    from kb_server.vector_store import VectorStore

    store = VectorStore()
    registry = IngestRegistry()

    await store.connect()
    registry.connect()

    if clean:
        log.info("Limpando KB e registry antes de reingerir...")
        await store.client.delete_collection(store.collection)
        await store._ensure_collection()
        registry.reset()
        log.info("KB e registry limpos.")

    if single_file:
        files = [single_file]
        docs_root = single_file.parent
    else:
        files = [
            f
            for f in docs_path.rglob("*")
            if f.is_file() and f.suffix.lower() in EXT_TYPE_MAP
        ]
        docs_root = docs_path

    # ── Detecta arquivos removidos do disco
    if sync and not single_file:
        _sync_deleted(registry, docs_root, files)

    total_chunks = 0
    total_ok = 0
    total_skipped = 0
    total_errors = 0
    start_time = time.time()

    sem = asyncio.Semaphore(workers)

    async def process_one(f: Path) -> tuple[int, str]:
        """Process a single file and return (chunk_count, status)."""
        async with sem:
            n, status = await process_file(
                f, docs_root, store, registry, product, force
            )
            return n, status

    results = await asyncio.gather(*[process_one(f) for f in files])

    for n, status in results:
        total_chunks += n
        if status == "ok":
            total_ok += 1
        elif status == "skipped":
            total_skipped += 1
        elif status == "error":
            total_errors += 1

    elapsed = time.time() - start_time
    summary = registry.summary()

    log.info(
        f"\n{'='*50}\n"
        f"Ingestão concluída em {elapsed:.1f}s\n"
        f"  Novos/atualizados:  {total_ok}\n"
        f"  Sem alterações:     {total_skipped}\n"
        f"  Erros:              {total_errors}\n"
        f"  Chunks indexados:   {total_chunks}\n"
        f"\n  Registry total:\n"
        f"  Documentos OK:      {summary['ok']}\n"
        f"  Com erro:           {summary['errors']}\n"
        f"  Total chunks KB:    {summary['total_chunks']}\n"
        f"{'='*50}"
    )

    if total_errors > 0:
        log.info("Arquivos com erro:")
        for e in registry.list_errors():
            log.info(f"  {e['path']}: {e['error_msg']}")

    registry.close()


def _sync_deleted(registry, docs_root: Path, current_files: list[Path]):
    """Marca no registry arquivos que foram removidos do disco."""
    existing_rel = {str(f.relative_to(docs_root)) for f in current_files}
    all_indexed = registry.list_all(status="ok")
    removed = 0
    for rec in all_indexed:
        if rec["path"] not in existing_rel:
            registry.mark_deleted(rec["path"])
            removed += 1
    if removed:
        log.info(
            f"  {removed} arquivo(s) removidos do disco "
            f"marcados como 'deleted'"
        )


def cmd_status(args):
    """Display registry status report without running ingestion."""
    import datetime

    from ingest.core.metadata import IngestRegistry

    registry = IngestRegistry()
    registry.connect()

    summary = registry.summary()

    def ts(t):
        """Format a Unix timestamp as a date string.

        Args:
            t: Unix timestamp (float).

        Returns:
            Formatted date string 'YYYY-MM-DD HH:MM', or em-dash if None.
        """
        if not t:
            return "—"
        return datetime.datetime.fromtimestamp(t).strftime("%Y-%m-%d %H:%M")

    print("\n" + "=" * 55)
    print(" KB RAG — Status do Registry")
    print("=" * 55)
    print(f"  Documentos indexados (OK):  {summary['ok'] or 0}")
    print(f"  Com erro:                   {summary['errors'] or 0}")
    print(f"  Removidos do disco:         {summary['deleted'] or 0}")
    print(f"  Total de chunks:            {summary['total_chunks'] or 0}")
    print(f"  Primeira ingestão:          {ts(summary['first_indexed'])}")
    print(f"  Última ingestão:            {ts(summary['last_indexed'])}")

    if args.errors:
        erros = registry.list_errors()
        if erros:
            print(f"\n  Arquivos com erro ({len(erros)}):")
            for e in erros:
                print(f"    [{ts(e['indexed_at'])}] {e['path']}")
                print(f"      → {e['error_msg']}")
        else:
            print("\n  Nenhum erro registrado.")

    if args.list:
        all_files = registry.list_all()
        print(f"\n  Todos os arquivos ({len(all_files)}):")
        for r in all_files:
            icon = {"ok": "✓", "error": "✗", "deleted": "⌀"}.get(
                r["status"], "?"
            )
            print(
                f"    {icon} [{r['file_type']:5}] {r['path']}  ("
                f"{r['chunks']} chunks, {ts(r['indexed_at'])}"
                ")"
            )

    print("=" * 55 + "\n")
    registry.close()


def main():
    """CLI entry point for the ingest pipeline.

    Parses arguments and routes to either the ingest subcommand or the
    status subcommand. Supports backward-compatible direct flags.
    """
    parser = argparse.ArgumentParser(
        description="KB RAG — Ingest Pipeline"
    )
    sub = parser.add_subparsers(dest="cmd")

    # ── ingest (default)
    p_ingest = sub.add_parser("ingest", help="Ingere documentos (padrão)")
    p_ingest.add_argument(
        "--docs", type=Path, help="Pasta raiz dos documentos"
    )
    p_ingest.add_argument(
        "--file", type=Path, help="Arquivo único para ingerir"
    )
    p_ingest.add_argument(
        "--product", type=str, help="Nome do produto (override)"
    )
    p_ingest.add_argument(
        "--workers", type=int, default=2, help="Workers paralelos"
    )
    p_ingest.add_argument(
        "--clean", action="store_true", help="Limpa KB e registry antes"
    )
    p_ingest.add_argument(
        "--force", action="store_true", help="Re-ingere mesmo sem mudanças"
    )
    p_ingest.add_argument(
        "--sync",
        action="store_true",
        help="Marca deletados arquivos removidos do disco",
    )

    # ── status
    p_status = sub.add_parser("status", help="Exibe status do registry")
    p_status.add_argument(
        "--errors", action="store_true", help="Mostra arquivos com erro"
    )
    p_status.add_argument(
        "--list", action="store_true", help="Lista todos os arquivos"
    )

    # Compatibilidade retroativa: aceita flags direto sem subcomando

    parser.add_argument("--docs", type=Path)
    parser.add_argument("--file", type=Path)
    parser.add_argument("--product", type=str)
    parser.add_argument("--workers", type=int, default=2)
    parser.add_argument("--clean", action="store_true")
    parser.add_argument("--force", action="store_true")
    parser.add_argument("--sync", action="store_true")
    parser.add_argument(
        "--status", action="store_true", help="Atalho para o subcomando status"
    )
    parser.add_argument("--errors", action="store_true")
    parser.add_argument("--list", action="store_true")

    args = parser.parse_args()

    # Roteamento
    if args.cmd == "status" or getattr(args, "status", False):
        cmd_status(args)
        return

    if args.cmd == "ingest":
        ingest_args = args
    else:
        ingest_args = args  # flags diretas

    if not ingest_args.docs and not ingest_args.file:
        parser.error("Informe --docs ou --file")

    docs_path = ingest_args.docs or (
        ingest_args.file.parent if ingest_args.file else Path(".")
    )
    asyncio.run(
        run_ingest(
            docs_path=docs_path,
            product=ingest_args.product,
            workers=ingest_args.workers,
            single_file=ingest_args.file,
            clean=ingest_args.clean,
            force=getattr(ingest_args, "force", False),
            sync=getattr(ingest_args, "sync", False),
        )
    )


if __name__ == "__main__":
    main()
