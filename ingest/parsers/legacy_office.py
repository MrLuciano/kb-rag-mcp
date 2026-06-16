"""
Legacy Office format extractors.

Supported formats:
- .doc  — via docx2txt (handles old binary Word via zip/xml fallback)
- .xls  — via xlrd (Excel 97-2003)
- .ppt  — via python-pptx (best-effort; binary .ppt unsupported)
- .odt, .ods, .odp — via odfpy (OpenDocument formats)
- .wpd  — text extraction only (WordPerfect; no reliable Python parser)

All functions return list[dict] with keys:
    text: str   — extracted text content
    page: int|str|None  — page/sheet number or name, None if not applicable
"""

import logging
from pathlib import Path

log = logging.getLogger("kb-ingest.parsers.legacy")


def extract_doc(path: Path) -> list[dict]:
    """
    Extract text from a .doc file.

    Uses docx2txt as primary (handles .docx-in-.doc-extension cases common
    in newer Office suites). Falls back to python-docx for truly docx files.
    """
    path = Path(path)
    if not path.exists():
        return []
    try:
        import docx2txt

        text = docx2txt.process(str(path))
        if text and text.strip():
            return [{"text": text.strip(), "page": None}]
    except Exception as e:
        log.debug(f"  docx2txt failed for {path.name}: {e}")

    # Fallback: try python-docx (works if file is actually .docx)
    try:
        from docx import Document

        doc = Document(str(path))
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        if paragraphs:
            return [{"text": "\n".join(paragraphs), "page": None}]
    except Exception as e:
        log.debug(f"  python-docx failed for {path.name}: {e}")

    log.warning(f"  Could not extract text from .doc file: {path.name}")
    return []


def extract_xls(path: Path) -> list[dict]:
    """Extract text from a legacy .xls (Excel 97-2003) file."""
    path = Path(path)
    if not path.exists():
        return []
    try:
        import xlrd

        wb = xlrd.open_workbook(str(path))
        results = []
        for sheet_name in wb.sheet_names():
            ws = wb.sheet_by_name(sheet_name)
            rows = []
            for row_idx in range(ws.nrows):
                row_values = [
                    str(ws.cell_value(row_idx, col)) for col in range(ws.ncols)
                ]
                row_text = "\t".join(v for v in row_values if v.strip())
                if row_text.strip():
                    rows.append(row_text)
            if rows:
                results.append(
                    {
                        "text": f"[Sheet: {sheet_name}]\n" + "\n".join(rows),
                        "page": sheet_name,
                    }
                )
        return results
    except ImportError:
        log.error("  Install xlrd: pip install xlrd")
        return []
    except Exception as e:
        log.error(f"  Error extracting XLS {path.name}: {e}")
        return []


def extract_ppt(path: Path) -> list[dict]:
    """
    Extract text from a .ppt file (PowerPoint 97-2003).

    python-pptx can sometimes open .ppt files saved in compatibility mode.
    Binary .ppt files will fail — logs a warning and returns empty.
    """
    path = Path(path)
    if not path.exists():
        return []
    try:
        from pptx import Presentation

        prs = Presentation(str(path))
        results = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                results.append({"text": "\n".join(texts), "page": i})
        return results
    except Exception as e:
        log.warning(
            f"  Could not extract .ppt {path.name} "
            f"(binary format not supported): {e}"
        )
        return []


def extract_odt(path: Path) -> list[dict]:
    """Extract text from OpenDocument Text (.odt) files."""
    path = Path(path)
    if not path.exists():
        return []
    try:
        from odf.opendocument import load
        from odf.text import P

        doc = load(str(path))
        paragraphs = []
        for p in doc.getElementsByType(P):
            text = "".join(
                node.data for node in p.childNodes if hasattr(node, "data")
            ).strip()
            if text:
                paragraphs.append(text)
        return (
            [{"text": "\n".join(paragraphs), "page": None}]
            if paragraphs
            else []
        )
    except ImportError:
        log.error("  Install odfpy: pip install odfpy")
        return []
    except Exception as e:
        log.error(f"  Error extracting ODT {path.name}: {e}")
        return []


def extract_ods(path: Path) -> list[dict]:
    """Extract text from OpenDocument Spreadsheet (.ods) files."""
    path = Path(path)
    if not path.exists():
        return []
    try:
        from odf.opendocument import load
        from odf.table import Table, TableCell, TableRow
        from odf.text import P

        doc = load(str(path))
        results = []
        for table in doc.getElementsByType(Table):
            sheet_name = table.getAttribute("name") or "Sheet"
            rows = []
            for row in table.getElementsByType(TableRow):
                cells = []
                for cell in row.getElementsByType(TableCell):
                    cell_text = "".join(
                        "".join(
                            node.data
                            for node in p.childNodes
                            if hasattr(node, "data")
                        )
                        for p in cell.getElementsByType(P)
                    ).strip()
                    cells.append(cell_text)
                row_text = "\t".join(cells).strip()
                if row_text:
                    rows.append(row_text)
            if rows:
                results.append(
                    {
                        "text": f"[Sheet: {sheet_name}]\n" + "\n".join(rows),
                        "page": sheet_name,
                    }
                )
        return results
    except ImportError:
        log.error("  Install odfpy: pip install odfpy")
        return []
    except Exception as e:
        log.error(f"  Error extracting ODS {path.name}: {e}")
        return []


def extract_odp(path: Path) -> list[dict]:
    """Extract text from OpenDocument Presentation (.odp) files."""
    path = Path(path)
    if not path.exists():
        return []
    try:
        from odf.draw import Page
        from odf.opendocument import load
        from odf.text import P

        doc = load(str(path))
        results = []
        for i, page in enumerate(doc.getElementsByType(Page), 1):
            paragraphs = []
            for p in page.getElementsByType(P):
                text = "".join(
                    node.data for node in p.childNodes if hasattr(node, "data")
                ).strip()
                if text:
                    paragraphs.append(text)
            if paragraphs:
                results.append({"text": "\n".join(paragraphs), "page": i})
        return results
    except ImportError:
        log.error("  Install odfpy: pip install odfpy")
        return []
    except Exception as e:
        log.error(f"  Error extracting ODP {path.name}: {e}")
        return []


def extract_wpd(path: Path) -> list[dict]:
    """
    Extract text from WordPerfect (.wpd) files.

    No reliable Python parser exists for binary .wpd. Attempts to read
    as latin-1 and strips non-printable characters. Quality is low
    but better than nothing for simple documents.
    """
    path = Path(path)
    if not path.exists():
        return []
    try:
        import re

        raw = path.read_bytes()
        text = raw.decode("latin-1", errors="replace")
        text = re.sub(r"[^\x20-\x7e\x80-\xff\n\r\t]", " ", text)
        text = re.sub(r" {3,}", " ", text)
        text = text.strip()
        if text:
            log.warning(
                f"  .wpd extracted with heuristic text strip "
                f"(low quality): {path.name}"
            )
            return [{"text": text, "page": None}]
    except Exception as e:
        log.error(f"  Error extracting WPD {path.name}: {e}")
    return []
